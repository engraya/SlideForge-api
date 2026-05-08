from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.exceptions import PresentationGenerationError
from src.schemas.presentation import PresentationTheme, SlideContent
from src.utils.logging import get_logger

logger = get_logger(__name__)

THEMES: dict[PresentationTheme, dict] = {
    PresentationTheme.PROFESSIONAL: {
        "bg_color": RGBColor(0, 51, 102),
        "title_color": RGBColor(255, 255, 255),
        "bullet_color": RGBColor(220, 235, 255),
        "title_font_size": Pt(36),
        "bullet_font_size_normal": Pt(20),
        "bullet_font_size_dense": Pt(16),
    },
    PresentationTheme.MINIMAL: {
        "bg_color": RGBColor(255, 255, 255),
        "title_color": RGBColor(30, 30, 30),
        "bullet_color": RGBColor(60, 60, 60),
        "title_font_size": Pt(34),
        "bullet_font_size_normal": Pt(20),
        "bullet_font_size_dense": Pt(16),
    },
    PresentationTheme.VIBRANT: {
        "bg_color": RGBColor(18, 18, 50),
        "title_color": RGBColor(100, 220, 255),
        "bullet_color": RGBColor(240, 240, 240),
        "title_font_size": Pt(36),
        "bullet_font_size_normal": Pt(20),
        "bullet_font_size_dense": Pt(16),
    },
}


class PresentationService:
    """
    Builds PowerPoint files from validated SlideContent objects.
    Fixes the runtime crash: the original create_pptx() accessed slide_info["title"]
    (dict syntax) but received list[list[str]] — a TypeError on every real call.
    This service uses slide_data.title and slide_data.bullets (Pydantic attributes).
    """

    def build(
        self,
        slides: list[SlideContent],
        output_path: Path,
        theme: PresentationTheme = PresentationTheme.PROFESSIONAL,
    ) -> None:
        theme_cfg = THEMES.get(theme, THEMES[PresentationTheme.PROFESSIONAL])

        try:
            prs = Presentation()
            cover_title = slides[0].title if slides else "Presentation"
            self._add_title_slide(prs, cover_title, theme_cfg)

            for slide_data in slides:
                self._add_content_slide(prs, slide_data, theme_cfg)

            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))

            logger.info(
                "Presentation saved",
                extra={"filename": output_path.name, "slides": len(slides), "theme": theme},
            )

        except PresentationGenerationError:
            raise
        except Exception as exc:
            raise PresentationGenerationError(
                detail=f"Failed to build presentation: {exc}",
                context={"output_path": str(output_path), "slide_count": len(slides)},
            ) from exc

    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        theme_cfg: dict,
    ) -> None:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme_cfg["bg_color"]

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            para = title_shape.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.bold = True
            para.font.size = Pt(44)
            para.font.color.rgb = theme_cfg["title_color"]

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: SlideContent,
        theme_cfg: dict,
    ) -> None:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme_cfg["bg_color"]

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
        title_tf = title_box.text_frame
        title_para = title_tf.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.bold = True
        title_para.font.size = theme_cfg["title_font_size"]
        title_para.font.color.rgb = theme_cfg["title_color"]

        bullet_count = len(slide_data.bullets)
        font_size = (
            theme_cfg["bullet_font_size_normal"]
            if bullet_count <= 5
            else theme_cfg["bullet_font_size_dense"]
        )

        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.0)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for idx, bullet_text in enumerate(slide_data.bullets):
            para = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            para.text = f"•  {bullet_text}"
            para.font.size = font_size
            para.font.color.rgb = theme_cfg["bullet_color"]
            para.space_after = Pt(8)
            if idx == 0:
                para.font.bold = True
